import pytesseract
import os
import fitz
import io
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
class PdfToText():
    def __init__(self,file_name):
        self.file_name = file_name
        self.pages = {}
    
    def extract_text(self):
        doc = fitz.open(self.file_name)
        total_pages = len(doc)
        for page_num,page in enumerate(doc,start = 1):
            txt = page.get_text()
            if len(txt.strip()) >150:
                self.pages[page_num] = txt
            else:
                pix = page.get_pixmap(matrix=fitz.Matrix(2,2))
                img_bytes = pix.tobytes("png")
                image = Image.open(io.BytesIO(img_bytes))
                ocr_txt = pytesseract.image_to_string(image)
                self.pages[page_num] = ocr_txt
            progress = page_num / total_pages

            bar_length = 30
            filled = int(progress * bar_length)

            bar = "█" * filled + "-" * (bar_length - filled)
            print(f"\r[{bar}] {progress*100:.1f}% ({page_num}/{total_pages})",end="")
        print()       
        doc.close()
        return self.pages

class DocxToText():
    def __init__(self,filename):
        self.filename = filename
        self.pages={}
    def extract_text(self):
        text = ""
        doc = Document(self.filename)
        for para in doc.paragraphs:
            text+=para.text +"\n"
        self.pages[1] = text
        return self.pages


class PptxToText():
    def __init__(self,filename):
        self.filename = filename
        self.pages = {}
    def extract_text(self):
        ppt = Presentation(self.filename)
        total_pages = len(ppt.slides)
        for page_num,slide in enumerate(ppt.slides,start=1):
            slide_txt = ""
            for shape in slide.shapes:
                if hasattr(shape,"text"):
                    slide_txt +=shape.text +"\n"
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    img_bytes = image.blob
                    img_type = shape.image.ext.lower()
                    if img_type in ["png","jpg","jpeg","tiff","bmp"]:
                    
                        try:
                            pil_img = Image.open(io.BytesIO(img_bytes))
                            ocr_text = pytesseract.image_to_string(pil_img,config="--psm 6")
                            if ocr_text.strip():
                                slide_txt += "\n" + ocr_text
                        except Exception as e:
                            print(f"\nOCR failed on slide {page_num}:{e}")
            self.pages[page_num] = slide_txt
            progress = page_num / total_pages

            bar_length = 30
            filled = int(progress * bar_length)
            bar = "█" * filled + "-" * (bar_length - filled)
            print(f"\r[{bar}] {progress*100:.1f}% ({page_num}/{total_pages})",end="")
        print() 
        return self.pages
    
class TxttoText():
    def __init__(self,filename):
        self.filename = filename
        self.pages = {}
    def extract_text(self):
        with open(self.filename,"r",encoding = "utf-8") as f:
            text = f.read()
        self.pages[1] = text
        return self.pages